import os
import tempfile
import zipfile
from typing import BinaryIO, List
import fitz  # PyMuPDF
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
import io

class PDFConverter:
    @staticmethod
    def to_images(pdf_file: BinaryIO) -> bytes:
        """Convert PDF to images using PyMuPDF and return as ZIP file."""
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_pdf:
            temp_pdf.write(pdf_file.read())
            temp_pdf_path = temp_pdf.name

        try:
            doc = fitz.open(temp_pdf_path)
            zip_buffer = io.BytesIO()
            with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
                for i, page in enumerate(doc):
                    pix = page.get_pixmap(dpi=200)
                    img_bytes = pix.tobytes("png")
                    zip_file.writestr(f'page_{i+1}.png', img_bytes)
            return zip_buffer.getvalue()
        finally:
            os.unlink(temp_pdf_path)

    @staticmethod
    def to_word(pdf_file: BinaryIO) -> bytes:
        """Convert PDF to Word document."""
        # Save PDF to temporary file
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_pdf:
            temp_pdf.write(pdf_file.read())
            temp_pdf_path = temp_pdf.name

        try:
            # Open PDF
            doc = fitz.open(temp_pdf_path)
            word_doc = Document()

            # Extract text from each page
            for page in doc:
                text = page.get_text()
                word_doc.add_paragraph(text)

            # Save to memory
            docx_buffer = io.BytesIO()
            word_doc.save(docx_buffer)
            return docx_buffer.getvalue()
        finally:
            os.unlink(temp_pdf_path)

    @staticmethod
    def to_excel(pdf_file: BinaryIO) -> bytes:
        """Convert PDF to Excel spreadsheet."""
        # Save PDF to temporary file
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_pdf:
            temp_pdf.write(pdf_file.read())
            temp_pdf_path = temp_pdf.name

        try:
            # Open PDF
            doc = fitz.open(temp_pdf_path)
            wb = Workbook()
            ws = wb.active
            ws.title = "PDF Content"

            # Extract text from each page
            for i, page in enumerate(doc):
                text = page.get_text()
                # Split text into rows and add to Excel
                rows = text.split('\n')
                for row_idx, row in enumerate(rows):
                    cells = row.split('\t')
                    for col_idx, cell in enumerate(cells):
                        ws.cell(row=row_idx + 1, column=col_idx + 1, value=cell)

            # Save to memory
            xlsx_buffer = io.BytesIO()
            wb.save(xlsx_buffer)
            return xlsx_buffer.getvalue()
        finally:
            os.unlink(temp_pdf_path)

    @staticmethod
    def to_powerpoint(pdf_file: BinaryIO) -> bytes:
        """Convert PDF to PowerPoint presentation."""
        # Save PDF to temporary file
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_pdf:
            temp_pdf.write(pdf_file.read())
            temp_pdf_path = temp_pdf.name

        try:
            # Open PDF
            doc = fitz.open(temp_pdf_path)
            prs = Presentation()

            # Convert each page to a slide
            for page in doc:
                # Create a new slide
                slide_layout = prs.slide_layouts[6]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)

                # Convert page to image
                pix = page.get_pixmap()
                img_data = pix.tobytes("png")

                # Add image to slide
                left = top = 0
                width = prs.slide_width
                height = prs.slide_height
                slide.shapes.add_picture(io.BytesIO(img_data), left, top, width, height)

            # Save to memory
            pptx_buffer = io.BytesIO()
            prs.save(pptx_buffer)
            return pptx_buffer.getvalue()
        finally:
            os.unlink(temp_pdf_path) 